from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta de cores ──────────────────────────────────────────
BG_DARK   = RGBColor(0x13, 0x08, 0x26)   # #130826
PINK      = RGBColor(0xE9, 0x1E, 0x8C)   # #E91E8C
VIOLET    = RGBColor(0x9C, 0x27, 0xB0)   # #9C27B0
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0xB0, 0x90, 0xCC)
GREEN     = RGBColor(0x00, 0xE5, 0xA0)
WARN      = RGBColor(0xFF, 0xB3, 0x00)
DANGER    = RGBColor(0xFF, 0x4D, 0x6D)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # layout em branco

def new_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=BG_DARK):
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def accent_bar(slide, color=PINK, top=Inches(0), height=Inches(0.06)):
    bar = slide.shapes.add_shape(1, 0, top, prs.slide_width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def add_text(slide, text, left, top, width, height,
             size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_pill(slide, text, left, top, width=Inches(2.6), height=Inches(0.45),
             bg_color=PINK, text_color=WHITE, size=14):
    box = slide.shapes.add_shape(9, left, top, width, height)  # rounded rect
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = True
    run.font.color.rgb = text_color

def add_card(slide, left, top, width, height, title, items,
             title_color=PINK, item_color=WHITE, bg_c=None):
    if bg_c is None:
        bg_c = RGBColor(0x22, 0x0A, 0x3E)
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg_c
    box.line.color.rgb = PINK
    box.line.width = Pt(1)

    # título do card
    add_text(slide, title,
             left + Inches(0.18), top + Inches(0.12),
             width - Inches(0.36), Inches(0.4),
             size=15, bold=True, color=title_color)

    # itens
    y = top + Inches(0.55)
    for item in items:
        add_text(slide, f"• {item}",
                 left + Inches(0.18), y,
                 width - Inches(0.36), Inches(0.35),
                 size=12, color=item_color)
        y += Inches(0.32)

# ════════════════════════════════════════════════════════════
# SLIDE 1 — Capa
# ════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
accent_bar(s, PINK,   top=Inches(0),    height=Inches(0.06))
accent_bar(s, VIOLET, top=Inches(7.44), height=Inches(0.06))

# Gradiente visual com retângulos
grad = s.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
grad.fill.solid()
grad.fill.fore_color.rgb = RGBColor(0x1A, 0x04, 0x38)
grad.line.fill.background()

add_text(s, "Mini-Meeting Dashboard",
         Inches(1), Inches(1.8), Inches(11.33), Inches(1.4),
         size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Sistema de gestão de eventos com check-in digital, QR Code e assinatura eletrônica",
         Inches(2), Inches(3.4), Inches(9.33), Inches(0.8),
         size=20, color=MUTED, align=PP_ALIGN.CENTER)

add_pill(s, "AstraZeneca  ·  2025/2026",
         Inches(4.67), Inches(4.4), Inches(4), Inches(0.5),
         bg_color=RGBColor(0x33, 0x0A, 0x5C), text_color=MUTED, size=14)

# ════════════════════════════════════════════════════════════
# SLIDE 2 — Visão Geral
# ════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
accent_bar(s, PINK, height=Inches(0.05))

add_text(s, "Visão Geral do Sistema",
         Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         size=30, bold=True, color=WHITE)
add_text(s, "O que é e para que serve",
         Inches(0.6), Inches(0.95), Inches(12), Inches(0.4),
         size=16, color=MUTED)

boxes = [
    ("📋 Gestão de Meetings",   ["Criação e edição de eventos", "Controle de status (ativo/encerrado)", "Visão admin × usuário"]),
    ("👥 Controle de Participantes", ["Inscrição via link público", "Importação em lote por CSV", "Assinatura digital"]),
    ("✅ Check-in Digital",     ["QR Code individual por participante", "Scanner via câmera", "Check-in manual na busca"]),
    ("🔐 Autenticação",         ["Login com JWT", "Perfis: Admin e Usuário", "Convite por e-mail (48h)"]),
]

cols = [Inches(0.4), Inches(3.6), Inches(6.8), Inches(10.0)]
for i, (title, items) in enumerate(boxes):
    add_card(s, cols[i], Inches(1.55), Inches(2.8), Inches(2.8),
             title, items)

# ════════════════════════════════════════════════════════════
# SLIDE 3 — Stack Tecnológico
# ════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
accent_bar(s, VIOLET, height=Inches(0.05))

add_text(s, "Stack Tecnológico",
         Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         size=30, bold=True, color=WHITE)

layers = [
    ("FRONTEND",  PINK,   ["React 19 + Vite 8", "React Router v7", "Axios + JWT interceptors", "react-signature-canvas", "qrcode.react  ·  html5-qrcode"]),
    ("BACKEND",   VIOLET, ["Node.js + Express 5", "MongoDB Atlas + Mongoose", "JWT (jsonwebtoken) + bcryptjs", "Nodemailer (SMTP)", "uuid  ·  dotenv"]),
    ("INFRA",     GREEN,  ["Frontend → Cloudflare Pages", "Backend → Render.com (auto-deploy)", "DB → MongoDB Atlas (cloud)", "CI/CD → push no branch main", "Domínio → azminimeeting.grfapps.com.br"]),
]

for i, (layer, color, items) in enumerate(layers):
    left = Inches(0.4 + i * 4.3)
    box = s.shapes.add_shape(1, left, Inches(1.3), Inches(4.0), Inches(5.6))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1E, 0x07, 0x38)
    box.line.color.rgb = color
    box.line.width = Pt(2)

    pill = s.shapes.add_shape(9, left + Inches(0.15), Inches(1.15), Inches(1.5), Inches(0.38))
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    tf = pill.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = layer
    r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = BG_DARK

    y = Inches(1.55)
    for item in items:
        add_text(s, f"  {item}",
                 left + Inches(0.2), y, Inches(3.6), Inches(0.38),
                 size=13, color=WHITE)
        y += Inches(0.44)

# ════════════════════════════════════════════════════════════
# SLIDE 4 — Arquitetura
# ════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
accent_bar(s, PINK, height=Inches(0.05))

add_text(s, "Arquitetura da Aplicação",
         Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         size=30, bold=True, color=WHITE)

# Diagrama simplificado
nodes = [
    (Inches(0.3),  Inches(3.0), Inches(2.4), Inches(1.2), "🌐 Browser\n(Cloudflare Pages)", PINK),
    (Inches(3.4),  Inches(1.5), Inches(2.6), Inches(1.2), "⚛️ React App\nCloudflare Pages", PINK),
    (Inches(3.4),  Inches(4.5), Inches(2.6), Inches(1.2), "📱 Mobile\n(link público)", MUTED),
    (Inches(7.2),  Inches(3.0), Inches(2.6), Inches(1.2), "🖥️ Express API\nRender.com", VIOLET),
    (Inches(10.4), Inches(1.5), Inches(2.4), Inches(1.0), "🍃 MongoDB\nAtlas", GREEN),
    (Inches(10.4), Inches(4.0), Inches(2.4), Inches(1.0), "📧 SMTP\nNodemailer", WARN),
]

for (l, t, w, h, txt, col) in nodes:
    box = s.shapes.add_shape(9, l, t, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x22, 0x08, 0x40)
    box.line.color.rgb = col
    box.line.width = Pt(2)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE

# setas (linhas)
arrows = [
    (Inches(2.7), Inches(3.6), Inches(3.4), Inches(2.1)),
    (Inches(2.7), Inches(3.6), Inches(3.4), Inches(5.1)),
    (Inches(6.0), Inches(2.1), Inches(7.2), Inches(3.6)),
    (Inches(6.0), Inches(5.1), Inches(7.2), Inches(3.6)),
    (Inches(9.8), Inches(3.6), Inches(10.4), Inches(2.0)),
    (Inches(9.8), Inches(3.6), Inches(10.4), Inches(4.5)),
]
for (x1,y1,x2,y2) in arrows:
    connector = s.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = MUTED
    connector.line.width = Pt(1.5)

add_text(s, "HTTPS / REST API",
         Inches(5.8), Inches(2.6), Inches(1.6), Inches(0.35),
         size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 5 — Perfis de Acesso
# ════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
accent_bar(s, VIOLET, height=Inches(0.05))

add_text(s, "Perfis de Acesso",
         Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         size=30, bold=True, color=WHITE)

# Admin
box = s.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.6))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1E, 0x07, 0x38)
box.line.color.rgb = PINK; box.line.width = Pt(2)

add_text(s, "🔑  ADMIN", Inches(0.7), Inches(1.25), Inches(5), Inches(0.5),
         size=22, bold=True, color=PINK)

admin_items = [
    "Vê todos os meetings de todos os usuários",
    "Edita e exclui qualquer meeting",
    "Gerencia usuários (ativar/desativar/remover)",
    "Envia convites por e-mail para novos usuários",
    "Acessa o scanner de QR Code",
    "Importa participantes via CSV",
]
y = Inches(1.85)
for item in admin_items:
    add_text(s, f"✓  {item}", Inches(0.8), y, Inches(5.2), Inches(0.38), size=13, color=WHITE)
    y += Inches(0.42)

# Usuário
box2 = s.shapes.add_shape(1, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.6))
box2.fill.solid(); box2.fill.fore_color.rgb = RGBColor(0x1E, 0x07, 0x38)
box2.line.color.rgb = VIOLET; box2.line.width = Pt(2)

add_text(s, "👤  USUÁRIO", Inches(7.2), Inches(1.25), Inches(5), Inches(0.5),
         size=22, bold=True, color=VIOLET)

user_items = [
    "Acessa apenas seus próprios meetings",
    "Pode ter 1 meeting ativo por vez",
    "Cria link de convite para participantes",
    "Importa participantes via CSV",
    "Acessa o scanner de QR Code",
    "Visualiza lista e status de check-ins",
]
y = Inches(1.85)
for item in user_items:
    add_text(s, f"✓  {item}", Inches(7.2), y, Inches(5.2), Inches(0.38), size=13, color=WHITE)
    y += Inches(0.42)

# ════════════════════════════════════════════════════════════
# SLIDE 6 — Fluxo do Participante
# ════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
accent_bar(s, GREEN, height=Inches(0.05))

add_text(s, "Fluxo do Participante",
         Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         size=30, bold=True, color=WHITE)

steps = [
    (GREEN,  "1", "Admin cria\no Meeting",       "Define título, local,\ndata e horário"),
    (PINK,   "2", "Admin importa\nCSV",           "Lista de participantes\ncom nome, email, CRM"),
    (VIOLET, "3", "Participante\nrecebe link",    "Link do QR Code\ncompartilhado pelo admin"),
    (WARN,   "4", "Busca seu\nQR Code",           "Digita nome ou e-mail\nna página de busca"),
    (GREEN,  "5", "Apresenta na\nentrada",        "QR Code escaneado\nou check-in manual"),
]

for i, (col, num, title, sub) in enumerate(steps):
    left = Inches(0.3 + i * 2.55)

    circle = s.shapes.add_shape(9, left + Inches(0.65), Inches(1.3), Inches(1.1), Inches(1.1))
    circle.fill.solid(); circle.fill.fore_color.rgb = col
    circle.line.fill.background()
    tf = circle.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = BG_DARK

    add_text(s, title, left, Inches(2.55), Inches(2.4), Inches(0.7),
             size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, sub, left, Inches(3.3), Inches(2.4), Inches(0.7),
             size=12, color=MUTED, align=PP_ALIGN.CENTER)

    if i < 4:
        arr = s.shapes.add_connector(1,
            left + Inches(2.4), Inches(1.85),
            left + Inches(2.55), Inches(1.85))
        arr.line.color.rgb = MUTED; arr.line.width = Pt(2)

# ════════════════════════════════════════════════════════════
# SLIDE 7 — Funcionalidades: CSV + QR Code
# ════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
accent_bar(s, WARN, height=Inches(0.05))

add_text(s, "Importação CSV & QR Codes",
         Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         size=30, bold=True, color=WHITE)

# CSV
add_text(s, "📂  Importação via CSV", Inches(0.5), Inches(1.15), Inches(6), Inches(0.5),
         size=20, bold=True, color=WARN)

csv_items = [
    "Botão \"Importar CSV\" na tela do meeting",
    "Suporta vírgula ( , ) e ponto-e-vírgula ( ; )",
    "Colunas: nome, email, crm, uf, telefone, cidade",
    "Duplicatas são ignoradas automaticamente",
    "Até 500 participantes por importação",
    "Botão \"Modelo\" baixa CSV de exemplo",
]
y = Inches(1.75)
for item in csv_items:
    add_text(s, f"• {item}", Inches(0.7), y, Inches(5.8), Inches(0.36), size=13, color=WHITE)
    y += Inches(0.38)

# Exemplo CSV
box = s.shapes.add_shape(1, Inches(0.5), Inches(5.0), Inches(5.8), Inches(1.6))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x0D, 0x02, 0x20)
box.line.color.rgb = WARN; box.line.width = Pt(1)
add_text(s, "nome,email,crm,uf,telefone,cidade\nJoão Silva,joao@email.com,123456,SP,(11) 99999-9999,São Paulo\nMaria Costa,maria@email.com,654321,RJ,(21) 98888-8888,Rio de Janeiro",
         Inches(0.7), Inches(5.1), Inches(5.4), Inches(1.4),
         size=11, color=GREEN, italic=True)

# QR Code
add_text(s, "📱  Página de QR Codes", Inches(7.0), Inches(1.15), Inches(6), Inches(0.5),
         size=20, bold=True, color=PINK)

qr_items = [
    "Link público por meeting (sem login)",
    "Participante busca por nome ou e-mail",
    "QR Code individual gerado na tela",
    "Botão de check-in manual disponível",
    "Status de check-in atualizado em tempo real",
]
y = Inches(1.75)
for item in qr_items:
    add_text(s, f"• {item}", Inches(7.2), y, Inches(5.8), Inches(0.36), size=13, color=WHITE)
    y += Inches(0.38)

# ════════════════════════════════════════════════════════════
# SLIDE 8 — Scanner de Check-in
# ════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
accent_bar(s, GREEN, height=Inches(0.05))

add_text(s, "Scanner de Check-in",
         Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         size=30, bold=True, color=WHITE)

add_text(s, "Acesso exclusivo: organizador do meeting ou admin",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
         size=15, color=MUTED, italic=True)

features = [
    (GREEN,  "📷",  "Câmera\nTraseira",    "Ativa a câmera do\ndispositivo para leitura"),
    (PINK,   "🔢",  "Contador\nAo Vivo",   "Mostra X / Total\npresentes em tempo real"),
    (VIOLET, "✅",  "Feedback\nImediato",  "Verde = confirmado\nRoxo = já registrado"),
    (WARN,   "⏱️",  "Cooldown\n2.5s",     "Evita duplo scan\npor engano"),
]

for i, (col, icon, title, sub) in enumerate(features):
    left = Inches(0.5 + i * 3.2)
    box = s.shapes.add_shape(9, left, Inches(1.6), Inches(2.8), Inches(4.0))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1E, 0x07, 0x38)
    box.line.color.rgb = col; box.line.width = Pt(2)

    add_text(s, icon, left, Inches(2.0), Inches(2.8), Inches(0.8),
             size=36, align=PP_ALIGN.CENTER)
    add_text(s, title, left, Inches(2.9), Inches(2.8), Inches(0.6),
             size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, sub, left, Inches(3.6), Inches(2.8), Inches(0.7),
             size=12, color=MUTED, align=PP_ALIGN.CENTER)

add_text(s, "Rota:  /meetings/:id/scan  (autenticada)",
         Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
         size=13, color=MUTED, italic=True)

# ════════════════════════════════════════════════════════════
# SLIDE 9 — API Routes
# ════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
accent_bar(s, VIOLET, height=Inches(0.05))

add_text(s, "Rotas da API",
         Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         size=30, bold=True, color=WHITE)

routes = [
    ("POST", "/api/auth/login",                      "Público",      MUTED),
    ("POST", "/api/auth/register",                   "Público",      MUTED),
    ("GET",  "/api/users",                           "Admin",        PINK),
    ("POST", "/api/users/invite",                    "Admin",        PINK),
    ("GET",  "/api/meetings",                        "Autenticado",  VIOLET),
    ("POST", "/api/meetings",                        "Autenticado",  VIOLET),
    ("PUT",  "/api/meetings/:id",                    "Autenticado",  VIOLET),
    ("POST", "/api/meetings/:id/attendees/bulk",     "Autenticado",  VIOLET),
    ("GET",  "/api/meetings/invite/:token",          "Público",      MUTED),
    ("GET",  "/api/meetings/invite/:token/lookup",   "Público",      MUTED),
    ("POST", "/api/meetings/invite/:token/register", "Público",      MUTED),
    ("POST", "/api/meetings/checkin/:token",         "Público",      GREEN),
]

method_colors = {"GET": GREEN, "POST": PINK, "PUT": WARN, "DELETE": DANGER}

cols_w = [Inches(0.9), Inches(4.8), Inches(2.0)]
headers = ["Método", "Rota", "Acesso"]
header_x = [Inches(0.4), Inches(1.4), Inches(6.3)]

for i, h in enumerate(headers):
    add_text(s, h, header_x[i], Inches(1.05), cols_w[i], Inches(0.35),
             size=12, bold=True, color=MUTED)

y = Inches(1.45)
for method, path, access, acc_color in routes:
    row_bg = s.shapes.add_shape(1, Inches(0.4), y - Inches(0.04),
                                 Inches(8.5), Inches(0.34))
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = RGBColor(0x1A, 0x06, 0x30)
    row_bg.line.fill.background()

    m_col = method_colors.get(method, WHITE)
    add_text(s, method, Inches(0.4), y, Inches(0.85), Inches(0.3),
             size=11, bold=True, color=m_col)
    add_text(s, path, Inches(1.35), y, Inches(4.8), Inches(0.3),
             size=11, color=WHITE)
    add_text(s, access, Inches(6.25), y, Inches(2.0), Inches(0.3),
             size=11, color=acc_color)
    y += Inches(0.38)

# ════════════════════════════════════════════════════════════
# SLIDE 10 — Deploy & Infraestrutura
# ════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
accent_bar(s, PINK, height=Inches(0.05))

add_text(s, "Deploy & Infraestrutura",
         Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         size=30, bold=True, color=WHITE)

infra = [
    (PINK,   "🌐", "Frontend",    "Cloudflare Pages",
     ["azminimeeting.grfapps.com.br", "Auto-deploy no push do branch main",
      "Variáveis: VITE_API_URL"]),
    (VIOLET, "🖥️", "Backend",     "Render.com",
     ["https://mini-meeting-astrazeneca.onrender.com", "Node.js 24 · auto-deploy",
      "Root Dir: backend  ·  Start: node server.js"]),
    (GREEN,  "🍃", "Banco de Dados", "MongoDB Atlas",
     ["Cluster gratuito (M0)", "Variável: MONGO_URI",
      "Mongoose como ODM"]),
    (WARN,   "📧", "E-mail",      "SMTP (Nodemailer)",
     ["Envio de convites para novos usuários", "Variáveis: EMAIL_HOST/PORT/USER/PASS",
      "Fallback: link exibido na tela"]),
]

for i, (col, icon, name, provider, items) in enumerate(infra):
    left = Inches(0.3 + i * 3.25)
    box = s.shapes.add_shape(1, left, Inches(1.2), Inches(3.0), Inches(5.6))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1E, 0x07, 0x38)
    box.line.color.rgb = col; box.line.width = Pt(2)

    add_text(s, icon, left, Inches(1.3), Inches(3.0), Inches(0.6),
             size=30, align=PP_ALIGN.CENTER)
    add_text(s, name, left, Inches(1.95), Inches(3.0), Inches(0.45),
             size=15, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, provider, left, Inches(2.42), Inches(3.0), Inches(0.35),
             size=12, color=MUTED, align=PP_ALIGN.CENTER)

    y = Inches(2.9)
    for item in items:
        add_text(s, f"• {item}", left + Inches(0.15), y, Inches(2.7), Inches(0.38),
                 size=11, color=WHITE)
        y += Inches(0.4)

# ════════════════════════════════════════════════════════════
# SLIDE 11 — Obrigado
# ════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
accent_bar(s, PINK,   top=Inches(0),    height=Inches(0.06))
accent_bar(s, VIOLET, top=Inches(7.44), height=Inches(0.06))

add_text(s, "Mini-Meeting Dashboard",
         Inches(1), Inches(2.0), Inches(11.33), Inches(1.0),
         size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Gestão inteligente de eventos corporativos",
         Inches(2), Inches(3.1), Inches(9.33), Inches(0.6),
         size=20, color=MUTED, align=PP_ALIGN.CENTER)

add_pill(s, "🌐  azminimeeting.grfapps.com.br",
         Inches(3.5), Inches(4.0), Inches(6.33), Inches(0.55),
         bg_color=RGBColor(0x33, 0x0A, 0x5C), text_color=PINK, size=15)

# ────────────────────────────────────────────────────────────
out = r"C:\ativacoes\mini-meeting\Mini-Meeting-Dashboard.pptx"
prs.save(out)
print(f"✅  Apresentação salva em: {out}")
