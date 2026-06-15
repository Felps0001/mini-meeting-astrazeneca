from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

SS   = r"C:\ativacoes\mini-meeting\screenshots"
OUT  = r"C:\ativacoes\mini-meeting\Mini-Meeting-Manual.pptx"

BG    = RGBColor(0x13, 0x08, 0x26)
PINK  = RGBColor(0xE9, 0x1E, 0x8C)
VIOLT = RGBColor(0x9C, 0x27, 0xB0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xB0, 0x90, 0xCC)
GREEN = RGBColor(0x00, 0xE5, 0xA0)
WARN  = RGBColor(0xFF, 0xB3, 0x00)
CARD  = RGBColor(0x1E, 0x07, 0x38)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background()
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = PINK; bar.line.fill.background()
    return s

def txt(s, text, l, t, w, h, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color

def img_card(s, img_path, left, top, iw, ih, caption, step_n=None, step_color=PINK):
    """Adiciona frame + imagem + caption."""
    # frame
    frame = s.shapes.add_shape(1, left - Inches(0.08), top - Inches(0.08),
                                iw + Inches(0.16), ih + Inches(0.16))
    frame.fill.solid(); frame.fill.fore_color.rgb = CARD
    frame.line.color.rgb = step_color; frame.line.width = Pt(1.5)

    # imagem
    s.shapes.add_picture(img_path, left, top, iw, ih)

    # caption
    txt(s, caption, left, top + ih + Inches(0.06), iw, Inches(0.35),
        size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # número do passo
    if step_n:
        badge = s.shapes.add_shape(9, left - Inches(0.15), top - Inches(0.15),
                                    Inches(0.38), Inches(0.38))
        badge.fill.solid(); badge.fill.fore_color.rgb = step_color
        badge.line.fill.background()
        tf = badge.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(step_n)
        r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = BG

def header(s, title, subtitle=None):
    txt(s, title, Inches(0.5), Inches(0.12), Inches(12), Inches(0.55),
        size=26, bold=True, color=WHITE)
    if subtitle:
        txt(s, subtitle, Inches(0.5), Inches(0.65), Inches(12), Inches(0.35),
            size=13, color=MUTED, italic=True)

# ══════════════════════════════════════════════════
# SLIDE 1 — Capa
# ══════════════════════════════════════════════════
s = slide()
s.shapes.add_picture(os.path.join(SS, "01_login.png"),
                     Inches(0), Inches(0), prs.slide_width, prs.slide_height)
# overlay escuro
ov = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
from pptx.dml.color import RGBColor as _C
ov.fill.solid()
ov.fill.fore_color.rgb = _C(0x13, 0x08, 0x26)
ov.line.fill.background()

txt(s, "Mini-Meeting Dashboard",
    Inches(1), Inches(2.2), Inches(11), Inches(1.2),
    size=50, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Manual de uso do sistema",
    Inches(2), Inches(3.5), Inches(9), Inches(0.6),
    size=22, color=MUTED, align=PP_ALIGN.CENTER)
pill = s.shapes.add_shape(9, Inches(4.5), Inches(4.3), Inches(4.3), Inches(0.55))
pill.fill.solid(); pill.fill.fore_color.rgb = PINK; pill.line.fill.background()
tf = pill.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "azminimeeting.grfapps.com.br"
r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE

# ══════════════════════════════════════════════════
# SLIDE 2 — Login
# ══════════════════════════════════════════════════
s = slide()
header(s, "1. Como fazer login", "Acesse o sistema pelo link e informe suas credenciais")

# esquerda: screenshot
img_card(s, os.path.join(SS, "01_login.png"),
         Inches(0.4), Inches(1.1), Inches(5.5), Inches(3.4),
         "Tela de login vazia", 1, PINK)

img_card(s, os.path.join(SS, "02_login_preenchido.png"),
         Inches(6.3), Inches(1.1), Inches(5.5), Inches(3.4),
         "Login preenchido — clique em Entrar", 2, PINK)

# caixas de instrução
steps_box = s.shapes.add_shape(1, Inches(0.4), Inches(4.9), Inches(11.4), Inches(2.2))
steps_box.fill.solid(); steps_box.fill.fore_color.rgb = CARD
steps_box.line.color.rgb = PINK; steps_box.line.width = Pt(1)

txt(s, "Como acessar:", Inches(0.7), Inches(5.0), Inches(11), Inches(0.4),
    size=13, bold=True, color=PINK)
instrucoes = [
    "① Acesse  https://azminimeeting.grfapps.com.br",
    "② Digite seu e-mail e senha nos campos indicados",
    "③ Clique em Entrar — você será redirecionado ao Dashboard",
    "⚠️  Novos usuários só entram por convite do admin (link enviado por e-mail com validade de 48h)",
]
y = Inches(5.45)
for i in instrucoes:
    txt(s, i, Inches(0.7), y, Inches(12), Inches(0.35), size=12, color=WHITE)
    y += Inches(0.38)

# ══════════════════════════════════════════════════
# SLIDE 3 — Dashboard
# ══════════════════════════════════════════════════
s = slide()
header(s, "2. Dashboard — Visão geral", "Resumo de todos os meetings e acesso rápido às funcionalidades")

img_card(s, os.path.join(SS, "03_dashboard.png"),
         Inches(0.4), Inches(1.1), Inches(8.0), Inches(4.5),
         "Dashboard principal com contadores", step_n=None, step_color=VIOLT)

# callouts laterais
callouts = [
    (GREEN,  "📋 Total de Meetings",   "Quantidade total cadastrada"),
    (PINK,   "🟢 Meetings Ativos",     "Eventos em andamento"),
    (WARN,   "👥 Participantes",       "Soma de todos os inscritos"),
    (VIOLT,  "+ Novo Meeting",         "Atalho para criar um evento"),
]
y = Inches(1.2)
for col, title, desc in callouts:
    box = s.shapes.add_shape(1, Inches(8.8), y, Inches(4.1), Inches(0.9))
    box.fill.solid(); box.fill.fore_color.rgb = CARD
    box.line.color.rgb = col; box.line.width = Pt(1.5)
    txt(s, title, Inches(9.0), y + Inches(0.08), Inches(3.8), Inches(0.35),
        size=13, bold=True, color=col)
    txt(s, desc, Inches(9.0), y + Inches(0.48), Inches(3.8), Inches(0.32),
        size=11, color=MUTED)
    y += Inches(1.05)

# ══════════════════════════════════════════════════
# SLIDE 4 — Lista de Meetings
# ══════════════════════════════════════════════════
s = slide()
header(s, "3. Lista de Meetings", "Veja, filtre e gerencie todos os eventos")

img_card(s, os.path.join(SS, "04_meetings.png"),
         Inches(0.4), Inches(1.1), Inches(8.0), Inches(4.5),
         "Tela de lista de meetings com filtros e ações")

callouts = [
    (GREEN,  "🔗  Copiar link",   "Copia link de inscrição pública"),
    (PINK,   "✏️  Editar",        "Edita título, local, data"),
    (WARN,   "🔒  Encerrar",      "Muda status para encerrado"),
    (MUTED,  "❌  Cancelar",      "Cancela o meeting"),
    (RGBColor(0xFF,0x4D,0x6D), "🗑️  Excluir", "Remove permanentemente"),
]
y = Inches(1.2)
for col, title, desc in callouts:
    box = s.shapes.add_shape(1, Inches(8.8), y, Inches(4.1), Inches(0.85))
    box.fill.solid(); box.fill.fore_color.rgb = CARD
    box.line.color.rgb = col; box.line.width = Pt(1.5)
    txt(s, title, Inches(9.0), y + Inches(0.08), Inches(3.8), Inches(0.35),
        size=13, bold=True, color=col)
    txt(s, desc, Inches(9.0), y + Inches(0.46), Inches(3.8), Inches(0.30),
        size=11, color=MUTED)
    y += Inches(0.98)

# ══════════════════════════════════════════════════
# SLIDE 5 — Criar Meeting
# ══════════════════════════════════════════════════
s = slide()
header(s, "4. Como criar um Meeting", "Clique em '+ Novo Meeting' na tela de Meetings ou Dashboard")

img_card(s, os.path.join(SS, "05_novo_meeting.png"),
         Inches(0.4), Inches(1.1), Inches(7.5), Inches(4.5),
         "Formulário de criação de meeting")

campos = [
    (PINK,  "Título *",             "Nome do evento  (obrigatório)"),
    (PINK,  "Descrição",            "Objetivo ou detalhes  (opcional)"),
    (PINK,  "Local *",              "Endereço ou sala  (obrigatório)"),
    (PINK,  "Data *",               "Data do evento  (obrigatório)"),
    (PINK,  "Horário de início *",  "Hora de início  (obrigatório)"),
    (VIOLT, "Horário de término",   "Hora de encerramento  (opcional)"),
]
y = Inches(1.2)
for col, campo, desc in campos:
    txt(s, f"• {campo}", Inches(8.2), y, Inches(4.8), Inches(0.28),
        size=13, bold=True, color=col)
    txt(s, f"  {desc}", Inches(8.2), y + Inches(0.28), Inches(4.8), Inches(0.28),
        size=11, color=MUTED)
    y += Inches(0.62)

txt(s, "💡  Um usuário pode ter apenas 1 meeting ativo por vez.",
    Inches(8.2), Inches(5.6), Inches(4.8), Inches(0.4),
    size=12, color=WARN, bold=True)

# ══════════════════════════════════════════════════
# SLIDE 6 — Detalhe do Meeting
# ══════════════════════════════════════════════════
s = slide()
header(s, "5. Detalhe do Meeting", "Clique no nome de um meeting para ver os detalhes e participantes")

img_card(s, os.path.join(SS, "06_meeting_detalhe.png"),
         Inches(0.4), Inches(1.1), Inches(8.0), Inches(5.0),
         "Informações do evento + lista de participantes")

notes = [
    (GREEN,  "🔗 Copiar link de convite",   "Link público para inscrição"),
    (VIOLT,  "📱 QR Codes",                 "Abre página de busca de QR"),
    (WARN,   "📷 Escanear QR",              "Abre scanner de check-in"),
    (PINK,   "📂 Importar CSV",             "Upload de lista de participantes"),
    (MUTED,  "📋 Modelo",                   "Baixa CSV de exemplo"),
    (RGBColor(0xFF,0x4D,0x6D), "🔒 / ❌",  "Encerrar ou cancelar"),
]
y = Inches(1.2)
for col, title, desc in notes:
    box = s.shapes.add_shape(1, Inches(8.7), y, Inches(4.2), Inches(0.82))
    box.fill.solid(); box.fill.fore_color.rgb = CARD
    box.line.color.rgb = col; box.line.width = Pt(1.5)
    txt(s, title, Inches(8.9), y + Inches(0.06), Inches(3.8), Inches(0.32),
        size=12, bold=True, color=col)
    txt(s, desc, Inches(8.9), y + Inches(0.44), Inches(3.8), Inches(0.28),
        size=11, color=MUTED)
    y += Inches(0.93)

# ══════════════════════════════════════════════════
# SLIDE 7 — Importar CSV
# ══════════════════════════════════════════════════
s = slide()
header(s, "6. Importar participantes via CSV",
       "Na tela do Meeting → clique em '📂 Importar CSV'")

img_card(s, os.path.join(SS, "07_meeting_participantes.png"),
         Inches(0.4), Inches(1.1), Inches(7.5), Inches(4.3),
         "Botões Importar CSV e Modelo na seção de participantes")

txt(s, "Formato do arquivo CSV:", Inches(8.2), Inches(1.2), Inches(4.8), Inches(0.4),
    size=14, bold=True, color=PINK)

code_box = s.shapes.add_shape(1, Inches(8.2), Inches(1.65), Inches(4.7), Inches(1.4))
code_box.fill.solid(); code_box.fill.fore_color.rgb = RGBColor(0x0D, 0x02, 0x20)
code_box.line.color.rgb = VIOLT; code_box.line.width = Pt(1)
txt(s, "nome,email,crm,uf,telefone,cidade\nJoão Silva,joao@email.com,\n123456,SP,(11)99999-9999,São Paulo",
    Inches(8.35), Inches(1.72), Inches(4.4), Inches(1.2),
    size=11, color=GREEN, italic=True)

regras = [
    "• Delimitador: vírgula ( , ) ou ponto-e-vírgula ( ; )",
    "• Apenas nome e e-mail são obrigatórios",
    "• Duplicatas (mesmo e-mail) são ignoradas",
    "• Máximo de 500 participantes por importação",
    "• Clique em '📋 Modelo' para baixar o template",
]
y = Inches(3.2)
for r in regras:
    txt(s, r, Inches(8.2), y, Inches(4.8), Inches(0.32), size=12, color=WHITE)
    y += Inches(0.38)

# ══════════════════════════════════════════════════
# SLIDE 8 — QR Code dos Participantes
# ══════════════════════════════════════════════════
s = slide()
header(s, "7. Página de QR Codes dos Participantes",
       "Na tela do Meeting → clique em '📱 QR Codes' (abre em nova aba)")

img_card(s, os.path.join(SS, "09_qrcode_busca.png"),
         Inches(0.4), Inches(1.1), Inches(5.5), Inches(3.8),
         "Página de busca — participante digita o nome")

img_card(s, os.path.join(SS, "10_qrcode_resultado.png"),
         Inches(6.2), Inches(1.1), Inches(5.5), Inches(3.8),
         "Resultado com QR Code e botão de check-in")

box = s.shapes.add_shape(1, Inches(0.4), Inches(5.1), Inches(11.4), Inches(1.9))
box.fill.solid(); box.fill.fore_color.rgb = CARD
box.line.color.rgb = VIOLT; box.line.width = Pt(1)

passos = [
    "1️⃣  Admin copia o link '📱 QR Codes' e compartilha (WhatsApp, e-mail, etc.)",
    "2️⃣  Participante acessa o link no próprio celular e digita seu nome ou e-mail",
    "3️⃣  Aparece o QR Code individual — apresenta na entrada para escaneamento",
    "4️⃣  Alternativa: clicar em '✅ Fazer Check-in' diretamente na página, sem scanner",
]
y = Inches(5.2)
for p in passos:
    txt(s, p, Inches(0.7), y, Inches(11), Inches(0.35), size=12, color=WHITE)
    y += Inches(0.38)

# ══════════════════════════════════════════════════
# SLIDE 9 — Scanner
# ══════════════════════════════════════════════════
s = slide()
header(s, "8. Scanner de Check-in",
       "Na tela do Meeting → clique em '📷 Escanear QR'  (somente organizador / admin)")

img_card(s, os.path.join(SS, "11_scanner.png"),
         Inches(0.4), Inches(1.1), Inches(7.5), Inches(4.8),
         "Scanner de QR Code com contador de presentes")

notas = [
    (GREEN,  "📷 Iniciar Scanner",   "Ativa câmera traseira do dispositivo"),
    (PINK,   "Contador ao vivo",     "Exibe X / Total presentes em tempo real"),
    (GREEN,  "✅ Check-in confirmado", "Nome aparece em verde"),
    (VIOLT,  "ℹ️ Já registrado",     "Aparece em roxo — não conta duplicado"),
    (RGBColor(0xFF,0x4D,0x6D), "❌ QR inválido", "Token desconhecido ou expirado"),
    (WARN,   "Cooldown 2,5s",        "Evita dupla leitura acidental"),
]
y = Inches(1.2)
for col, title, desc in notas:
    box = s.shapes.add_shape(1, Inches(8.2), y, Inches(4.7), Inches(0.78))
    box.fill.solid(); box.fill.fore_color.rgb = CARD
    box.line.color.rgb = col; box.line.width = Pt(1.5)
    txt(s, title, Inches(8.4), y + Inches(0.06), Inches(4.3), Inches(0.3),
        size=12, bold=True, color=col)
    txt(s, desc, Inches(8.4), y + Inches(0.42), Inches(4.3), Inches(0.28),
        size=11, color=MUTED)
    y += Inches(0.88)

# ══════════════════════════════════════════════════
# SLIDE 10 — Gerenciar Usuários
# ══════════════════════════════════════════════════
s = slide()
header(s, "9. Gerenciar Usuários  (somente Admin)",
       "Menu 'Usuários' na navbar — convide novos usuários e gerencie os existentes")

img_card(s, os.path.join(SS, "08_admin_users.png"),
         Inches(0.4), Inches(1.1), Inches(7.8), Inches(4.8),
         "Tela de gerenciamento de usuários")

acoes = [
    (PINK,   "🔗 Gerar link",        "Gera convite único com validade de 48h"),
    (VIOLT,  "📥 Importar CSV",      "Cria várias contas de uma vez"),
    (MUTED,  "📄 Baixar template",   "Modelo CSV para importação em massa"),
    (WARN,   "🚫 Desativar",         "Bloqueia acesso do usuário sem excluir"),
    (RGBColor(0xFF,0x4D,0x6D), "🗑️ Excluir", "Remove o usuário permanentemente"),
]
y = Inches(1.2)
for col, title, desc in acoes:
    box = s.shapes.add_shape(1, Inches(8.5), y, Inches(4.4), Inches(0.82))
    box.fill.solid(); box.fill.fore_color.rgb = CARD
    box.line.color.rgb = col; box.line.width = Pt(1.5)
    txt(s, title, Inches(8.7), y + Inches(0.06), Inches(4.0), Inches(0.32),
        size=13, bold=True, color=col)
    txt(s, desc, Inches(8.7), y + Inches(0.44), Inches(4.0), Inches(0.28),
        size=11, color=MUTED)
    y += Inches(0.93)

# ══════════════════════════════════════════════════
# SLIDE 11 — Resumo / Encerramento
# ══════════════════════════════════════════════════
s = slide()

bg2 = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg2.fill.solid(); bg2.fill.fore_color.rgb = RGBColor(0x13, 0x08, 0x26)
bg2.line.fill.background()

txt(s, "Resumo das funcionalidades",
    Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
    size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

features = [
    (PINK,   "🔐 Login",           "E-mail + senha com JWT"),
    (VIOLT,  "📋 Meetings",        "Criar, editar, encerrar"),
    (GREEN,  "👥 Participantes",   "Manual, CSV ou link público"),
    (WARN,   "📂 Importar CSV",    "Até 500 por importação"),
    (PINK,   "📱 QR Codes",        "Busca por nome/e-mail"),
    (GREEN,  "✅ Check-in",        "Scanner ou botão manual"),
    (VIOLT,  "📷 Scanner",         "Câmera com contador ao vivo"),
    (MUTED,  "👑 Admin",           "Gestão completa de usuários"),
]

cols = [Inches(0.3), Inches(3.55), Inches(6.8), Inches(10.05)]
rows = [Inches(0.95), Inches(2.55), Inches(4.15)]

for idx, (col, title, desc) in enumerate(features):
    c = idx % 4
    r = idx // 4
    left = cols[c]
    top  = rows[r]
    box = s.shapes.add_shape(9, left, top, Inches(2.9), Inches(1.35))
    box.fill.solid(); box.fill.fore_color.rgb = CARD
    box.line.color.rgb = col; box.line.width = Pt(2)
    txt(s, title, left, top + Inches(0.1), Inches(2.9), Inches(0.45),
        size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, desc, left, top + Inches(0.6), Inches(2.9), Inches(0.45),
        size=12, color=MUTED, align=PP_ALIGN.CENTER)

pill = s.shapes.add_shape(9, Inches(3.5), Inches(6.85), Inches(6.3), Inches(0.5))
pill.fill.solid(); pill.fill.fore_color.rgb = PINK; pill.line.fill.background()
tf = pill.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "🌐  azminimeeting.grfapps.com.br"
r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE

# ─────────────────────────────────────────────────
prs.save(OUT)
print(f"✅  Apresentação salva em: {OUT}")
