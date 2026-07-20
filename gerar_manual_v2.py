from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

SS  = r"C:\ativacoes\mini-meeting\screenshots"
OUT = r"C:\ativacoes\mini-meeting\Mini-Meeting-Manual-v2.pptx"

BG    = RGBColor(0x13, 0x08, 0x26)
PINK  = RGBColor(0xE9, 0x1E, 0x8C)
VIOLT = RGBColor(0x9C, 0x27, 0xB0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xB0, 0x90, 0xCC)
GREEN = RGBColor(0x00, 0xE5, 0xA0)
WARN  = RGBColor(0xFF, 0xB3, 0x00)
RED   = RGBColor(0xFF, 0x4D, 0x6D)
CARD  = RGBColor(0x1E, 0x07, 0x38)
CODE  = RGBColor(0x0D, 0x02, 0x20)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── helpers ──────────────────────────────────────────────────
def new_slide(bar_color=PINK):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background()
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = bar_color; bar.line.fill.background()
    return s

def txt(s, text, l, t, w, h, size=13, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color

def header(s, title, sub=None, bar=PINK):
    txt(s, title, Inches(0.5), Inches(0.1), Inches(12.3), Inches(0.55),
        size=26, bold=True, color=WHITE)
    if sub:
        txt(s, sub, Inches(0.5), Inches(0.63), Inches(12.3), Inches(0.35),
            size=12, color=MUTED, italic=True)

def pic(s, fname, l, t, w, h, caption=None, num=None, col=PINK):
    """Insere imagem com borda, número e legenda opcionais."""
    frame = s.shapes.add_shape(1, l - Inches(0.07), t - Inches(0.07),
                                w + Inches(0.14), h + Inches(0.14))
    frame.fill.solid(); frame.fill.fore_color.rgb = CARD
    frame.line.color.rgb = col; frame.line.width = Pt(1.5)
    s.shapes.add_picture(os.path.join(SS, fname), l, t, w, h)
    if caption:
        txt(s, caption, l, t + h + Inches(0.07), w, Inches(0.32),
            size=10, color=MUTED, align=PP_ALIGN.CENTER)
    if num:
        b = s.shapes.add_shape(9, l - Inches(0.14), t - Inches(0.14),
                                Inches(0.36), Inches(0.36))
        b.fill.solid(); b.fill.fore_color.rgb = col; b.line.fill.background()
        tf = b.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(num)
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = BG

def callout_list(s, left, top, items, width=Inches(4.5)):
    """Lista de callouts (cor, título, descrição)."""
    y = top
    for col, title, desc in items:
        box = s.shapes.add_shape(1, left, y, width, Inches(0.78))
        box.fill.solid(); box.fill.fore_color.rgb = CARD
        box.line.color.rgb = col; box.line.width = Pt(1.5)
        txt(s, title, left + Inches(0.18), y + Inches(0.06),
            width - Inches(0.22), Inches(0.3), size=12, bold=True, color=col)
        txt(s, desc,  left + Inches(0.18), y + Inches(0.42),
            width - Inches(0.22), Inches(0.28), size=10, color=MUTED)
        y += Inches(0.88)

def infobox(s, left, top, width, height, lines, title=None, col=PINK):
    box = s.shapes.add_shape(1, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = CARD
    box.line.color.rgb = col; box.line.width = Pt(1)
    y = top + Inches(0.1)
    if title:
        txt(s, title, left + Inches(0.15), y, width - Inches(0.2), Inches(0.35),
            size=13, bold=True, color=col)
        y += Inches(0.4)
    for line in lines:
        txt(s, line, left + Inches(0.15), y, width - Inches(0.2), Inches(0.32),
            size=11, color=WHITE)
        y += Inches(0.35)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ══════════════════════════════════════════════════════════════
s = new_slide(PINK)
s.shapes.add_picture(os.path.join(SS, "01_login.png"),
                     Inches(0), Inches(0), prs.slide_width, prs.slide_height)
ov = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
ov.fill.solid(); ov.fill.fore_color.rgb = BG; ov.line.fill.background()

txt(s, "Mini-Meeting Dashboard",
    Inches(1.2), Inches(2.0), Inches(10.9), Inches(1.2),
    size=50, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Manual completo do sistema",
    Inches(2), Inches(3.3), Inches(9.3), Inches(0.6),
    size=20, color=MUTED, align=PP_ALIGN.CENTER)

pill = s.shapes.add_shape(9, Inches(4.2), Inches(4.1), Inches(4.9), Inches(0.55))
pill.fill.solid(); pill.fill.fore_color.rgb = PINK; pill.line.fill.background()
tf = pill.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "🌐  azminimeeting.grfapps.com.br"
r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — LOGIN
# ══════════════════════════════════════════════════════════════
s = new_slide(PINK)
header(s, "1. Login", "Acesse o sistema informando seu e-mail e senha")

pic(s, "01_login.png", Inches(0.4), Inches(1.1), Inches(5.5), Inches(3.4),
    "Tela de login", 1, PINK)
pic(s, "02_login_preenchido.png", Inches(6.2), Inches(1.1), Inches(5.5), Inches(3.4),
    "Login preenchido — clique em Entrar", 2, PINK)

infobox(s, Inches(0.4), Inches(4.75), Inches(11.45), Inches(2.35), [
    "① Acesse  https://azminimeeting.grfapps.com.br",
    "② Preencha seu e-mail e senha",
    "③ Clique em Entrar — você será redirecionado ao Dashboard",
    "⚠️  Novos usuários só acessam via convite do admin (link por e-mail, válido por 48h)",
], title="Como acessar:", col=PINK)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — DASHBOARD
# ══════════════════════════════════════════════════════════════
s = new_slide(VIOLT)
header(s, "2. Dashboard", "Visão geral de todos os meetings e atalhos rápidos")

pic(s, "03_dashboard.png", Inches(0.4), Inches(1.1), Inches(8.1), Inches(4.6),
    "Dashboard principal")

callout_list(s, Inches(8.8), Inches(1.1), [
    (GREEN, "📋  Total de Meetings",  "Quantidade total cadastrada no sistema"),
    (PINK,  "🟢  Meetings Ativos",    "Eventos em andamento"),
    (WARN,  "👥  Participantes",      "Soma de todos os inscritos"),
    (MUTED, "Meetings Recentes",      "Lista dos últimos eventos criados"),
    (VIOLT, "+ Novo Meeting",         "Atalho para criação rápida"),
], width=Inches(4.1))

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — LISTA DE MEETINGS
# ══════════════════════════════════════════════════════════════
s = new_slide(VIOLT)
header(s, "3. Lista de Meetings", "Menu 'Meetings' na navbar — veja, filtre e gerencie")

pic(s, "04_meetings.png", Inches(0.4), Inches(1.1), Inches(8.1), Inches(4.6),
    "Tela de meetings com filtros de status")

callout_list(s, Inches(8.8), Inches(1.1), [
    (GREEN, "🔗  Link de convite",    "Copia link público de inscrição"),
    (PINK,  "✏️  Editar",             "Edita título, local, data e horário"),
    (WARN,  "🔒  Encerrar",           "Muda status para encerrado"),
    (MUTED, "❌  Cancelar",           "Cancela o meeting"),
    (RED,   "🗑️  Excluir",            "Remove permanentemente"),
], width=Inches(4.1))

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — CRIAR MEETING
# ══════════════════════════════════════════════════════════════
s = new_slide(GREEN)
header(s, "4. Criar novo Meeting", "Clique em '+ Novo Meeting' no Dashboard ou na lista de meetings")

pic(s, "05_novo_meeting.png", Inches(0.4), Inches(1.1), Inches(7.6), Inches(4.5),
    "Formulário de criação de meeting")

callout_list(s, Inches(8.3), Inches(1.1), [
    (PINK,  "Título *",              "Nome do evento  (obrigatório)"),
    (PINK,  "Descrição",             "Objetivo/detalhes  (opcional)"),
    (PINK,  "Local *",               "Sala ou endereço  (obrigatório)"),
    (PINK,  "Data *",                "Data do evento  (obrigatório)"),
    (PINK,  "Horário de início *",   "Hora de início  (obrigatório)"),
    (VIOLT, "Horário de término",    "Hora de encerramento  (opcional)"),
], width=Inches(4.6))

infobox(s, Inches(8.3), Inches(6.6), Inches(4.6), Inches(0.6), [
    "💡  Usuário: apenas 1 meeting ativo por vez",
], col=WARN)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — DETALHE DO MEETING (TOPO)
# ══════════════════════════════════════════════════════════════
s = new_slide(PINK)
header(s, "5. Detalhe do Meeting — Botões de ação",
       "Clique no título de qualquer meeting para abrir esta tela")

pic(s, "06_meeting_detalhe_topo.png", Inches(0.4), Inches(1.1), Inches(8.1), Inches(5.0),
    "Tela de detalhe com informações e botões de ação")

callout_list(s, Inches(8.8), Inches(1.1), [
    (GREEN, "🔗  Copiar link de convite", "Link público para inscrição de participantes"),
    (VIOLT, "📱  QR Codes",              "Abre busca de QR Code (nova aba)"),
    (WARN,  "📷  Escanear QR",           "Abre scanner de check-in"),
    (PINK,  "✏️  Editar",                "Edita informações do meeting"),
    (RED,   "🗑️  Excluir",               "Remove o meeting (somente admin)"),
    (MUTED, "🔒 Encerrar / ❌ Cancelar", "Muda o status do evento"),
], width=Inches(4.1))

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — DETALHE DO MEETING (PARTICIPANTES + CSV)
# ══════════════════════════════════════════════════════════════
s = new_slide(PINK)
header(s, "6. Participantes e Importação CSV",
       "Role a tela do meeting para baixo para ver a lista de participantes")

pic(s, "07_meeting_participantes.png", Inches(0.4), Inches(1.1), Inches(8.1), Inches(4.6),
    "Seção de participantes com botões de importação e atualização")

callout_list(s, Inches(8.8), Inches(1.1), [
    (GREEN, "📂  Importar CSV",      "Abre seletor de arquivo .csv"),
    (VIOLT, "📋  Modelo",            "Baixa um CSV de exemplo preenchido"),
    (MUTED, "🔄  Atualizar",         "Recarrega a lista sem sair da página"),
    (RED,   "✕  Remover",            "Cancela inscrição de um participante"),
], width=Inches(4.1))

# código CSV
code = s.shapes.add_shape(1, Inches(8.8), Inches(4.7), Inches(4.1), Inches(2.45))
code.fill.solid(); code.fill.fore_color.rgb = CODE
code.line.color.rgb = VIOLT; code.line.width = Pt(1)
txt(s, "Formato do CSV:", Inches(9.0), Inches(4.78), Inches(3.7), Inches(0.3),
    size=11, bold=True, color=VIOLT)
txt(s, "nome,email,crm,uf,telefone,cidade\nJoão Silva,joao@email.com,123456,\nSP,(11) 99999-9999,São Paulo",
    Inches(9.0), Inches(5.12), Inches(3.7), Inches(1.0),
    size=10, color=GREEN, italic=True)
txt(s, "• Duplicatas (mesmo e-mail) são ignoradas\n• Máx. 500 participantes por importação",
    Inches(9.0), Inches(6.15), Inches(3.7), Inches(0.7),
    size=10, color=MUTED)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — QR CODES DOS PARTICIPANTES
# ══════════════════════════════════════════════════════════════
s = new_slide(VIOLT)
header(s, "7. Página de QR Codes",
       "Clique em '📱 QR Codes' na tela do meeting — abre em nova aba")

pic(s, "09_qrcode_vazia.png", Inches(0.4), Inches(1.1), Inches(5.5), Inches(3.8),
    "Página aberta — participante digita o nome", 1, VIOLT)
pic(s, "10_qrcode_resultado.png", Inches(6.1), Inches(1.1), Inches(5.5), Inches(3.8),
    "QR Code encontrado + botão de check-in", 2, VIOLT)

infobox(s, Inches(0.4), Inches(5.15), Inches(11.45), Inches(2.0), [
    "1️⃣  Admin copia o link '📱 QR Codes' na tela do meeting e compartilha (WhatsApp, e-mail…)",
    "2️⃣  Participante acessa no celular, digita nome ou e-mail e encontra seu QR Code",
    "3️⃣  Apresenta o QR Code para ser escaneado  OU  clica em '✅ Fazer Check-in' manualmente",
], title="Fluxo de uso:", col=VIOLT)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — SCANNER DE CHECK-IN
# ══════════════════════════════════════════════════════════════
s = new_slide(GREEN)
header(s, "8. Scanner de Check-in",
       "Clique em '📷 Escanear QR' na tela do meeting  —  somente organizador ou admin")

pic(s, "11_scanner.png", Inches(0.4), Inches(1.1), Inches(7.8), Inches(5.0),
    "Tela do scanner com contador de presentes e feedback visual")

callout_list(s, Inches(8.6), Inches(1.1), [
    (GREEN, "📷  Iniciar Scanner",      "Ativa câmera traseira do dispositivo"),
    (PINK,  "Contador ao vivo",         "Exibe X / Total presentes em tempo real"),
    (GREEN, "✅  Check-in confirmado",  "Card verde com nome do participante"),
    (VIOLT, "ℹ️  Já registrado",        "Card roxo — sem duplicar no contador"),
    (RED,   "❌  QR inválido",          "Token desconhecido ou outro meeting"),
    (WARN,  "Cooldown 2,5s",            "Evita dupla leitura acidental"),
], width=Inches(4.3))

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — GERENCIAR USUÁRIOS (ADMIN)
# ══════════════════════════════════════════════════════════════
s = new_slide(WARN)
header(s, "9. Gerenciar Usuários  (somente Admin)",
       "Menu 'Usuários' na navbar — convide, importe e gerencie usuários do sistema")

pic(s, "08_admin_users.png", Inches(0.4), Inches(1.1), Inches(7.8), Inches(4.8),
    "Tela de gerenciamento de usuários")

callout_list(s, Inches(8.6), Inches(1.1), [
    (PINK,  "🔗  Gerar link",          "Convite único com validade de 48h"),
    (VIOLT, "📥  Importar CSV",        "Cria várias contas de uma vez"),
    (MUTED, "📄  Baixar template",     "Modelo CSV para importação em massa"),
    (WARN,  "🚫  Desativar",           "Bloqueia sem excluir o usuário"),
    (RED,   "🗑️  Excluir",             "Remove permanentemente"),
], width=Inches(4.3))

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — MAPA DO SISTEMA
# ══════════════════════════════════════════════════════════════
s = new_slide(PINK)
header(s, "10. Mapa de funcionalidades", "Visão completa de todas as features do sistema")

blocks = [
    # (col, emoji, título, subtítulo, col_idx, row_idx)
    (PINK,  "🔐", "Login",             "JWT · 7 dias",                    0, 0),
    (VIOLT, "📋", "Meetings",          "Criar · Editar · Encerrar",       1, 0),
    (GREEN, "👥", "Participantes",     "Manual · CSV · Link público",     2, 0),
    (WARN,  "📂", "Importar CSV",      "Até 500 · Dedup automático",      3, 0),
    (PINK,  "📱", "QR Codes",          "Busca por nome/e-mail",           0, 1),
    (GREEN, "✅", "Check-in manual",   "Botão direto na busca",           1, 1),
    (VIOLT, "📷", "Scanner QR",        "Câmera · Contador ao vivo",       2, 1),
    (WARN,  "👑", "Admin — Usuários",  "Convite · CSV · Ativar/Remover",  3, 1),
    (MUTED, "📊", "Dashboard",         "Contadores em tempo real",        0, 2),
    (GREEN, "🔔", "Notificações",      "Toast de sucesso/erro/aviso",     1, 2),
    (PINK,  "⏳", "Loading global",    "Barra de progresso nas requests",  2, 2),
    (VIOLT, "🛡️", "Segurança JWT",     "Expiry check · 401 auto-logout",  3, 2),
]

col_x = [Inches(0.3), Inches(3.55), Inches(6.8), Inches(10.05)]
row_y = [Inches(1.05), Inches(2.65), Inches(4.25)]

for col, icon, title, sub, ci, ri in blocks:
    bx = s.shapes.add_shape(9, col_x[ci], row_y[ri], Inches(2.9), Inches(1.35))
    bx.fill.solid(); bx.fill.fore_color.rgb = CARD
    bx.line.color.rgb = col; bx.line.width = Pt(2)
    txt(s, icon + "  " + title, col_x[ci], row_y[ri] + Inches(0.1),
        Inches(2.9), Inches(0.45), size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, sub, col_x[ci], row_y[ri] + Inches(0.62),
        Inches(2.9), Inches(0.45), size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — ENCERRAMENTO
# ══════════════════════════════════════════════════════════════
s = new_slide(PINK)

txt(s, "Mini-Meeting Dashboard",
    Inches(1), Inches(2.2), Inches(11.33), Inches(1.1),
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Sistema completo de gestão de eventos com check-in digital",
    Inches(1.5), Inches(3.4), Inches(10.33), Inches(0.6),
    size=18, color=MUTED, align=PP_ALIGN.CENTER)

pill = s.shapes.add_shape(9, Inches(4.2), Inches(4.2), Inches(4.9), Inches(0.55))
pill.fill.solid(); pill.fill.fore_color.rgb = PINK; pill.line.fill.background()
tf = pill.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "🌐  azminimeeting.grfapps.com.br"
r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE

# ─────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Salvo em: {OUT}")
